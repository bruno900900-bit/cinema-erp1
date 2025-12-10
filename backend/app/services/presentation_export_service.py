"""
Export Service - Geração de apresentações PowerPoint
"""
from typing import List, Optional
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from sqlalchemy.orm import Session, joinedload
import httpx

from ..models.location import Location
from ..models.location_photo import LocationPhoto
from ..schemas.presentation_export import PresentationExportRequest


class PresentationExportService:
    """Serviço para exportar locações para PowerPoint"""

    def __init__(self, db: Session):
        self.db = db
        # Cores do tema escuro
        self.colors = {
            'background': RGBColor(0x1A, 0x1A, 0x2E),  # Dark blue-gray
            'primary': RGBColor(0x4F, 0x46, 0xE5),     # Indigo
            'accent': RGBColor(0x06, 0xB6, 0xD4),       # Cyan
            'text_primary': RGBColor(0xFF, 0xFF, 0xFF), # White
            'text_secondary': RGBColor(0x9C, 0xA3, 0xAF), # Gray
            'card_bg': RGBColor(0x1E, 0x29, 0x3B),     # Darker blue
        }

    def create_presentation(self, request: PresentationExportRequest) -> bytes:
        """Cria apresentação PowerPoint das locações"""

        # Reordenar IDs conforme order
        ordered_ids = [request.location_ids[i] for i in request.order]

        # Buscar locações
        locations = self.db.query(Location).options(
            joinedload(Location.photos)
        ).filter(Location.id.in_(ordered_ids)).all()

        # Ordenar conforme a ordem recebida
        location_map = {loc.id: loc for loc in locations}
        ordered_locations = [location_map[lid] for lid in ordered_ids if lid in location_map]

        # Mapear fotos selecionadas
        selected_photos_map = {}
        if request.selected_photos:
            for sel in request.selected_photos:
                selected_photos_map[sel.location_id] = sel.photo_ids

        # Criar apresentação
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Slide de capa
        title = request.title or "Apresentação de Locações"
        subtitle = request.subtitle
        self._add_cover_slide(prs, title, subtitle)

        # Slides de locações
        for location in ordered_locations:
            # Filtrar fotos se seleção específica
            photos = []
            if request.include_photos and location.photos:
                if location.id in selected_photos_map:
                    photo_ids = selected_photos_map[location.id]
                    photos = [p for p in location.photos if p.id in photo_ids]
                else:
                    photos = list(location.photos)

            self._add_location_slide(prs, location, photos, request.include_summary)

        # Slide de resumo se solicitado
        if request.include_summary:
            self._add_summary_slide(prs, ordered_locations)

        # Slide final
        self._add_closing_slide(prs)

        # Salvar em buffer e retornar bytes
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()

    def _add_cover_slide(self, prs: Presentation, title: str, subtitle: Optional[str]):
        """Adiciona slide de capa"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Background escuro
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['background']
        background.line.fill.background()

        # Título
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.33), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(48)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.colors['text_primary']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Subtítulo
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2),
                Inches(11.33), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            tf.paragraphs[0].text = subtitle
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.color.rgb = self.colors['text_secondary']
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Linha decorativa
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.66), Inches(3.9),
            Inches(4), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors['primary']
        line.line.fill.background()

    def _add_location_slide(self, prs: Presentation, location: Location, photos: List, include_info: bool):
        """Adiciona slide de locação"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['background']
        background.line.fill.background()

        # Header com título
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['card_bg']
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.35),
            Inches(8), Inches(0.6)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = location.title or f"Locação #{location.id}"
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.colors['text_primary']

        # Cidade/Estado
        city_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.85),
            Inches(8), Inches(0.4)
        )
        tf = city_box.text_frame
        tf.paragraphs[0].text = f"{location.city or ''}, {location.state or ''}"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = self.colors['text_secondary']

        content_top = 1.5

        # Área de fotos (se disponível)
        if photos:
            photo = photos[0]
            photo_url = photo.url if hasattr(photo, 'url') and photo.url else None

            if photo_url:
                try:
                    # Baixar imagem
                    with httpx.Client(timeout=10.0) as client:
                        response = client.get(photo_url)
                        if response.status_code == 200:
                            image_bytes = BytesIO(response.content)
                            # Adicionar imagem grande
                            slide.shapes.add_picture(
                                image_bytes,
                                Inches(0.5), Inches(content_top),
                                width=Inches(7)
                            )
                except Exception as e:
                    print(f"Erro ao baixar imagem: {e}")

            # Miniaturas adicionais
            if len(photos) > 1:
                thumb_left = 7.8
                thumb_top = content_top
                thumb_size = 1.5
                for i, p in enumerate(photos[1:5]):  # Max 4 thumbnails
                    p_url = p.url if hasattr(p, 'url') and p.url else None
                    if not p_url:
                        continue
                    try:
                        with httpx.Client(timeout=10.0) as client:
                            resp = client.get(p_url)
                            if resp.status_code == 200:
                                img_bytes = BytesIO(resp.content)
                                slide.shapes.add_picture(
                                    img_bytes,
                                    Inches(thumb_left), Inches(thumb_top + i * (thumb_size + 0.2)),
                                    width=Inches(thumb_size)
                                )
                    except Exception:
                        pass

        # Informações da locação
        info_left = 8 if photos else 0.5
        info_width = 4.8 if photos else 12

        if include_info:
            info_box = slide.shapes.add_textbox(
                Inches(info_left), Inches(content_top),
                Inches(info_width), Inches(5)
            )
            tf = info_box.text_frame
            tf.word_wrap = True

            # Adicionar informações
            infos = []
            sector = getattr(location, 'sector_type', None)
            if sector:
                sector_val = sector.value if hasattr(sector, 'value') else str(sector)
                infos.append(f"Setor: {sector_val}")
            space = getattr(location, 'space_type', None)
            if space:
                space_val = space.value if hasattr(space, 'value') else str(space)
                infos.append(f"Tipo: {space_val}")
            if location.capacity:
                infos.append(f"Capacidade: {location.capacity} pessoas")
            daily_rate = getattr(location, 'daily_rate', None) or getattr(location, 'price_day_cinema', None)
            if daily_rate:
                infos.append(f"Diária: R$ {daily_rate:,.2f}")
            if location.description:
                desc = location.description[:200] + "..." if len(location.description or "") > 200 else location.description
                infos.append(f"\n{desc}")

            for idx, info in enumerate(infos):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = info
                p.font.size = Pt(12)
                p.font.color.rgb = self.colors['text_secondary']
                p.space_after = Pt(8)

    def _add_summary_slide(self, prs: Presentation, locations: List[Location]):
        """Adiciona slide de resumo"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['background']
        background.line.fill.background()

        # Título
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = "Resumo das Locações"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.colors['text_primary']

        # Stats
        total_locations = len(locations)
        total_capacity = sum(loc.capacity or 0 for loc in locations)
        total_daily = sum(loc.daily_rate or 0 for loc in locations)

        stats_text = f"""
Total de Locações: {total_locations}
Capacidade Total: {total_capacity} pessoas
Valor Total Diárias: R$ {total_daily:,.2f}
        """.strip()

        stats_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(6), Inches(2)
        )
        tf = stats_box.text_frame
        tf.paragraphs[0].text = stats_text
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = self.colors['text_secondary']

    def _add_closing_slide(self, prs: Presentation):
        """Adiciona slide de encerramento"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['background']
        background.line.fill.background()

        # Mensagem
        msg_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(11.33), Inches(1)
        )
        tf = msg_box.text_frame
        tf.paragraphs[0].text = "Obrigado!"
        tf.paragraphs[0].font.size = Pt(48)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.colors['text_primary']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Rodapé
        footer_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5),
            Inches(11.33), Inches(0.5)
        )
        tf = footer_box.text_frame
        tf.paragraphs[0].text = "Cinema ERP - Sistema de Gestão de Locações"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = self.colors['text_secondary']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
